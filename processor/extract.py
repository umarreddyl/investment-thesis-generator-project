from pptx import Presentation
import sys,json

def extract_text(file_path):
    ppt=Presentation(file_path)
    data=[]
    for slide in ppt.slides:
        text=[]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text.append(para.text)
        data.append(" ".join(text))
    return data

if __name__=="__main__":
    file_path=sys.argv[1]
    result=extract_text(file_path)
    print(json.dumps(result))